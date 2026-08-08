import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def create_edumind_deck():
    prs = Presentation()
    # Set 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette Definitions
    BG_DARK = RGBColor(15, 23, 42)      # Deep Navy #0F172A
    CARD_BG = RGBColor(30, 41, 59)      # Slate Dark #1E293B
    CARD_BORDER = RGBColor(51, 65, 85)  # Border #334155
    TEXT_MAIN = RGBColor(248, 250, 252) # White #F8FAFC
    TEXT_MUTED = RGBColor(148, 163, 184) # Muted Gray #94A3B8
    ACCENT_INDIGO = RGBColor(99, 102, 241) # Indigo #6366F1
    ACCENT_GREEN = RGBColor(16, 185, 129) # Emerald Green #10B981
    ACCENT_BLUE = RGBColor(59, 130, 246)  # Blue #3B82F6
    ACCENT_PURPLE = RGBColor(168, 85, 247) # Purple #A855F7

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_DARK

    def add_header(slide, title_text, category_text="EDUMIND • EDGEMINDS INTERNSHIP 2026"):
        # Header banner text box
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p1 = tf.paragraphs[0]
        p1.text = category_text.upper()
        p1.font.size = Pt(10)
        p1.font.bold = True
        p1.font.color.rgb = ACCENT_INDIGO

        p2 = tf.add_paragraph()
        p2.text = title_text
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_MAIN
        p2.space_before = Pt(4)

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
        return shape

    # =========================================================================
    # SLIDE 1: Cover Slide
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)

    # Top Banner Card
    banner = add_card(slide1, Inches(0.8), Inches(0.6), Inches(11.733), Inches(0.6), bg_color=RGBColor(30, 27, 75), border_color=ACCENT_INDIGO)
    tf_banner = banner.text_frame
    tf_banner.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_b = tf_banner.paragraphs[0]
    p_b.text = "🏆 EDGEMINDS INTERNSHIP 2026 – PROJECT SUMMARY AND FINAL REPORT"
    p_b.font.size = Pt(12)
    p_b.font.bold = True
    p_b.font.color.rgb = ACCENT_INDIGO
    p_b.alignment = PP_ALIGN.CENTER

    # Main Project Title Box
    title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(2.2))
    tf_t = title_box.text_frame
    tf_t.word_wrap = True

    p_t1 = tf_t.paragraphs[0]
    p_t1.text = "EduMind"
    p_t1.font.size = Pt(46)
    p_t1.font.bold = True
    p_t1.font.color.rgb = ACCENT_GREEN

    p_t2 = tf_t.add_paragraph()
    p_t2.text = "Edge-AI Offline RAG MCQ & Pagewise Note Generator Engine"
    p_t2.font.size = Pt(22)
    p_t2.font.bold = True
    p_t2.font.color.rgb = TEXT_MAIN
    p_t2.space_before = Pt(8)

    p_t3 = tf_t.add_paragraph()
    p_t3.text = "100% Local Intelligence Powered by Llama 3.2 1B & NVIDIA Jetson Platform"
    p_t3.font.size = Pt(14)
    p_t3.font.color.rgb = TEXT_MUTED
    p_t3.space_before = Pt(6)

    # Project Details Cards (2 Columns)
    card_left = add_card(slide1, Inches(0.8), Inches(4.0), Inches(5.7), Inches(2.8))
    tf_cl = card_left.text_frame
    tf_cl.margin_left = tf_cl.margin_top = tf_cl.margin_right = tf_cl.margin_bottom = Inches(0.3)
    tf_cl.word_wrap = True

    p_l1 = tf_cl.paragraphs[0]
    p_l1.text = "👥 Team Information"
    p_l1.font.size = Pt(16)
    p_l1.font.bold = True
    p_l1.font.color.rgb = ACCENT_BLUE

    items_l = [
        ("Team Name:", "Edge Thinkers"),
        ("Team Leader:", "Aditya Verma"),
        ("Track:", "Edge AI / Local RAG & SLM Integration"),
    ]
    for label, val in items_l:
        p = tf_cl.add_paragraph()
        p.space_before = Pt(10)
        run1 = p.add_run()
        run1.text = f"{label} "
        run1.font.bold = True
        run1.font.size = Pt(13)
        run1.font.color.rgb = TEXT_MUTED
        run2 = p.add_run()
        run2.text = val
        run2.font.bold = True
        run2.font.size = Pt(13)
        run2.font.color.rgb = TEXT_MAIN

    card_right = add_card(slide1, Inches(6.833), Inches(4.0), Inches(5.7), Inches(2.8))
    tf_cr = card_right.text_frame
    tf_cr.margin_left = tf_cr.margin_top = tf_cr.margin_right = tf_cr.margin_bottom = Inches(0.3)
    tf_cr.word_wrap = True

    p_r1 = tf_cr.paragraphs[0]
    p_r1.text = "🏫 Academic Affiliation"
    p_r1.font.size = Pt(16)
    p_r1.font.bold = True
    p_r1.font.color.rgb = ACCENT_PURPLE

    items_r = [
        ("Institution:", "Madan Mohan Malaviya University of Technology (MMMUT)"),
        ("Department:", "Electronics & Communication Engineering (ECE)"),
        ("Deployment Target:", "NVIDIA Jetson Orin / L4T CUDA Acceleration"),
    ]
    for label, val in items_r:
        p = tf_cr.add_paragraph()
        p.space_before = Pt(10)
        run1 = p.add_run()
        run1.text = f"{label} "
        run1.font.bold = True
        run1.font.size = Pt(13)
        run1.font.color.rgb = TEXT_MUTED
        run2 = p.add_run()
        run2.text = val
        run2.font.bold = True
        run2.font.size = Pt(13)
        run2.font.color.rgb = TEXT_MAIN

    # =========================================================================
    # SLIDE 2: Meet the Team
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "Meet the Team — Edge Thinkers")

    team = [
        ("Aditya Verma", "Lead Developer (Team Leader)", "RAG Architecture, Local LLM Integration & Quality Pipeline", ACCENT_INDIGO),
        ("Satyam Kumar Mishra", "Systems Engineer", "Streamlit Interface & Jetson Hardware Deployment", ACCENT_GREEN),
        ("Sachin Maurya", "ML / Data Pipeline Engineer", "PDF Ingestion, Regex OCR Cleaning & Chunking", ACCENT_BLUE),
        ("Ratandeep Shukla", "Testing & Evaluation Lead", "Quiz Quality Gates & Benchmark Evaluation", ACCENT_PURPLE),
    ]

    card_width = Inches(2.7)
    card_gap = Inches(0.3)
    left_start = Inches(0.8)

    for i, (name, role, contrib, color) in enumerate(team):
        c_left = left_start + i * (card_width + card_gap)
        c_card = add_card(slide2, c_left, Inches(1.6), card_width, Inches(5.2))
        tf_c = c_card.text_frame
        tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = Inches(0.25)
        tf_c.word_wrap = True

        # Photo Placeholder Circle Box
        circle = slide2.shapes.add_shape(MSO_SHAPE.OVAL, c_left + Inches(0.75), Inches(2.0), Inches(1.2), Inches(1.2))
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(15, 23, 42)
        circle.line.color.rgb = color
        circle.line.width = Pt(2)

        tf_circ = circle.text_frame
        tf_circ.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_circ = tf_circ.paragraphs[0]
        p_circ.text = name[0] + (name.split()[1][0] if len(name.split()) > 1 else "")
        p_circ.font.size = Pt(20)
        p_circ.font.bold = True
        p_circ.font.color.rgb = color
        p_circ.alignment = PP_ALIGN.CENTER

        # Member Info
        tb_info = slide2.shapes.add_textbox(c_left + Inches(0.1), Inches(3.4), card_width - Inches(0.2), Inches(3.2))
        tf_i = tb_info.text_frame
        tf_i.word_wrap = True

        p_n = tf_i.paragraphs[0]
        p_n.text = name
        p_n.font.size = Pt(16)
        p_n.font.bold = True
        p_n.font.color.rgb = TEXT_MAIN
        p_n.alignment = PP_ALIGN.CENTER

        p_r = tf_i.add_paragraph()
        p_r.text = role
        p_r.font.size = Pt(12)
        p_r.font.bold = True
        p_r.font.color.rgb = color
        p_r.alignment = PP_ALIGN.CENTER
        p_r.space_before = Pt(4)

        p_div = tf_i.add_paragraph()
        p_div.text = "───────────────"
        p_div.font.size = Pt(8)
        p_div.font.color.rgb = CARD_BORDER
        p_div.alignment = PP_ALIGN.CENTER
        p_div.space_before = Pt(6)

        p_c = tf_i.add_paragraph()
        p_c.text = contrib
        p_c.font.size = Pt(11)
        p_c.font.color.rgb = TEXT_MUTED
        p_c.alignment = PP_ALIGN.CENTER
        p_c.space_before = Pt(8)

    # =========================================================================
    # SLIDE 3: Problem Statement
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "Problem Statement — Educational AI Limitations")

    # Main Problem Box
    prob_card = add_card(slide3, Inches(0.8), Inches(1.5), Inches(11.733), Inches(1.6), bg_color=RGBColor(30, 27, 75), border_color=RGBColor(239, 68, 68))
    tf_p = prob_card.text_frame
    tf_p.margin_left = tf_p.margin_top = tf_p.margin_right = tf_p.margin_bottom = Inches(0.25)
    tf_p.word_wrap = True

    p_p1 = tf_p.paragraphs[0]
    p_p1.text = "⚠️ The Problem:"
    p_p1.font.size = Pt(16)
    p_p1.font.bold = True
    p_p1.font.color.rgb = RGBColor(248, 113, 113)

    p_p2 = tf_p.add_paragraph()
    p_p2.text = "Traditional AI learning tools rely heavily on cloud APIs, requiring continuous internet connectivity, incurring high recurring costs, and producing degraded assessments (verbatim fill-in-the-blank questions and broken PDF encodings)."
    p_p2.font.size = Pt(13)
    p_p2.font.color.rgb = TEXT_MAIN
    p_p2.space_before = Pt(4)

    # Target Audience Card
    aud_card = add_card(slide3, Inches(0.8), Inches(3.3), Inches(11.733), Inches(0.9))
    tf_a = aud_card.text_frame
    tf_a.margin_left = tf_a.margin_top = tf_a.margin_right = tf_a.margin_bottom = Inches(0.2)
    tf_a.word_wrap = True

    p_a1 = tf_a.paragraphs[0]
    p_a1.text = "👤 Who Faces It: University students, educators, and institutions in low-bandwidth, offline, or air-gapped study environments."
    p_a1.font.size = Pt(13)
    p_a1.font.bold = True
    p_a1.font.color.rgb = ACCENT_GREEN

    # 3 Pain Points Column Cards
    pain_points = [
        ("1. High Cloud Dependency", "High latency, expensive token costs, and total failure when internet connectivity is lost or throttled.", ACCENT_BLUE),
        ("2. Corrupted PDF Text", "Broken font encodings in raw lecture notes (e.g. 'c ntr l' instead of 'control') ruin vector retrieval accuracy.", ACCENT_PURPLE),
        ("3. Low MCQ Quality", "LLMs default to trivial fill-in-the-blank questions (rote recall) instead of testing deep conceptual understanding.", RGBColor(239, 68, 68)),
    ]

    p_width = Inches(3.7)
    p_gap = Inches(0.3)
    for i, (title, desc, color) in enumerate(pain_points):
        card = add_card(slide3, Inches(0.8) + i * (p_width + p_gap), Inches(4.4), p_width, Inches(2.4))
        tf_pt = card.text_frame
        tf_pt.margin_left = tf_pt.margin_top = tf_pt.margin_right = tf_pt.margin_bottom = Inches(0.2)
        tf_pt.word_wrap = True

        p_t = tf_pt.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = color

        p_d = tf_pt.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_before = Pt(8)

    # =========================================================================
    # SLIDE 4: Why This Problem Matters
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "Why This Problem Matters — Strategic Value")

    quads = [
        ("💡 1. Broad Educational Impact", "Enables uninterrupted, 100% offline self-assessment and revision for millions of students without expensive subscription cloud APIs.", ACCENT_INDIGO),
        ("🚀 2. Why Now?", "Recent breakthroughs in 1B Small Language Models (SLMs like Llama 3.2 1B) make real-time edge intelligence viable on compact embedded devices.", ACCENT_GREEN),
        ("🔒 3. Why Edge AI?", "Zero internet requirement, 100% data privacy for confidential university lecture notes, and sub-2s inference latency.", ACCENT_BLUE),
        ("🎯 4. Who Benefits?", "University students preparing for competitive exams, remote academic faculties, and air-gapped institutional study labs.", ACCENT_PURPLE),
    ]

    q_width = Inches(5.7)
    q_height = Inches(2.4)

    coords = [
        (Inches(0.8), Inches(1.6)),
        (Inches(6.833), Inches(1.6)),
        (Inches(0.8), Inches(4.3)),
        (Inches(6.833), Inches(4.3)),
    ]

    for i, (title, desc, color) in enumerate(quads):
        x, y = coords[i]
        q_card = add_card(slide4, x, y, q_width, q_height)
        tf_q = q_card.text_frame
        tf_q.margin_left = tf_q.margin_top = tf_q.margin_right = tf_q.margin_bottom = Inches(0.25)
        tf_q.word_wrap = True

        p_qt = tf_q.paragraphs[0]
        p_qt.text = title
        p_qt.font.size = Pt(16)
        p_qt.font.bold = True
        p_qt.font.color.rgb = color

        p_qd = tf_q.add_paragraph()
        p_qd.text = desc
        p_qd.font.size = Pt(12)
        p_qd.font.color.rgb = TEXT_MAIN
        p_qd.space_before = Pt(8)

    # =========================================================================
    # SLIDE 5: Our Solution
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "Our Solution — EduMind Edge-AI Engine")

    # One Line Solution Card
    sol_banner = add_card(slide5, Inches(0.8), Inches(1.5), Inches(11.733), Inches(1.1), bg_color=RGBColor(6, 78, 59), border_color=ACCENT_GREEN)
    tf_sb = sol_banner.text_frame
    tf_sb.margin_left = tf_sb.margin_top = tf_sb.margin_right = tf_sb.margin_bottom = Inches(0.2)
    tf_sb.word_wrap = True

    p_sb1 = tf_sb.paragraphs[0]
    p_sb1.text = "✨ Solution in One Line:"
    p_sb1.font.size = Pt(12)
    p_sb1.font.bold = True
    p_sb1.font.color.rgb = ACCENT_GREEN

    p_sb2 = tf_sb.add_paragraph()
    p_sb2.text = "An end-to-end, 100% offline Edge-AI RAG platform that transforms raw syllabus PDFs into validated conceptual MCQs and pagewise study notes."
    p_sb2.font.size = Pt(14)
    p_sb2.font.bold = True
    p_sb2.font.color.rgb = TEXT_MAIN
    p_sb2.space_before = Pt(2)

    # What We Built & Key Features (2 Columns)
    col1_card = add_card(slide5, Inches(0.8), Inches(2.8), Inches(5.7), Inches(2.6))
    tf_c1 = col1_card.text_frame
    tf_c1.margin_left = tf_c1.margin_top = tf_c1.margin_right = tf_c1.margin_bottom = Inches(0.25)
    tf_c1.word_wrap = True

    p_c1_t = tf_c1.paragraphs[0]
    p_c1_t.text = "⚙️ What We Built"
    p_c1_t.font.size = Pt(16)
    p_c1_t.font.bold = True
    p_c1_t.font.color.rgb = ACCENT_BLUE

    p_c1_d = tf_c1.add_paragraph()
    p_c1_d.text = "A lightweight Streamlit app running locally via Ollama (llama3.2:1b & nomic-embed-text) with automated regex OCR text repair and Python-level quality validation gates."
    p_c1_d.font.size = Pt(12)
    p_c1_d.font.color.rgb = TEXT_MAIN
    p_c1_d.space_before = Pt(8)

    col2_card = add_card(slide5, Inches(6.833), Inches(2.8), Inches(5.7), Inches(2.6))
    tf_c2 = col2_card.text_frame
    tf_c2.margin_left = tf_c2.margin_top = tf_c2.margin_right = tf_c2.margin_bottom = Inches(0.25)
    tf_c2.word_wrap = True

    p_c2_t = tf_c2.paragraphs[0]
    p_c2_t.text = "🌟 Key Features"
    p_c2_t.font.size = Pt(16)
    p_c2_t.font.bold = True
    p_c2_t.font.color.rgb = ACCENT_PURPLE

    feats = [
        "1. Sequential Pagewise Summarization (0% topic loss).",
        "2. Direct Conceptual MCQs (No fill-in-the-blank underscores).",
        "3. Python-Level Option Randomization & Page Citation.",
    ]
    for f in feats:
        p = tf_c2.add_paragraph()
        p.text = f
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_MAIN
        p.space_before = Pt(6)

    # Journey Pipeline Box at bottom
    j_card = add_card(slide5, Inches(0.8), Inches(5.6), Inches(11.733), Inches(1.2), bg_color=CARD_BG, border_color=ACCENT_INDIGO)
    tf_j = j_card.text_frame
    tf_j.margin_left = tf_j.margin_top = tf_j.margin_right = tf_j.margin_bottom = Inches(0.2)
    tf_j.word_wrap = True

    p_jt = tf_j.paragraphs[0]
    p_jt.text = "🔄 Input → Processing → Output Journey:"
    p_jt.font.size = Pt(12)
    p_jt.font.bold = True
    p_jt.font.color.rgb = ACCENT_INDIGO

    p_jd = tf_j.add_paragraph()
    p_jd.text = "Raw PDF Notes ➔ OCR Repair & Chunking ➔ Local Vector / Pagewise Extraction ➔ Ollama Inference & Quality Gate ➔ Exam-Ready Quiz & Revision Notes"
    p_jd.font.size = Pt(12)
    p_jd.font.bold = True
    p_jd.font.color.rgb = TEXT_MAIN
    p_jd.space_before = Pt(4)

    # =========================================================================
    # SLIDE 6: System Architecture
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "System Architecture — End-to-End Local Pipeline")

    steps_arch = [
        ("1. Input Layer", "Raw PDF / Word Notes\n(Multi-Page Upload)", ACCENT_BLUE),
        ("2. App / UI", "Streamlit Interface\n(Port 8505 / 8506)", ACCENT_INDIGO),
        ("3. Data Pipeline", "Regex OCR Repair &\nFAISS Vector Store", ACCENT_PURPLE),
        ("4. SLM Engine", "Ollama REST API\n(llama3.2:1b & nomic)", ACCENT_GREEN),
        ("5. Jetson Target", "NVIDIA Jetson Board\n(CUDA Acceleration)", RGBColor(245, 158, 11)),
        ("6. Validated Output", "Conceptual MCQs &\nPagewise Notes", RGBColor(236, 72, 153)),
    ]

    box_w = Inches(1.75)
    box_gap = Inches(0.2)
    start_x = Inches(0.8)

    for i, (title, desc, color) in enumerate(steps_arch):
        x = start_x + i * (box_w + box_gap)
        card = add_card(slide6, x, Inches(2.0), box_w, Inches(3.4))
        tf = card.text_frame
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.15)
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = color
        p_t.alignment = PP_ALIGN.CENTER

        p_div = tf.add_paragraph()
        p_div.text = "─────────"
        p_div.font.size = Pt(8)
        p_div.font.color.rgb = CARD_BORDER
        p_div.alignment = PP_ALIGN.CENTER

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_MAIN
        p_d.alignment = PP_ALIGN.CENTER
        p_d.space_before = Pt(8)

        # Arrow indicator between boxes
        if i < 5:
            arr_box = slide6.shapes.add_textbox(x + box_w, Inches(3.4), box_gap, Inches(0.6))
            tf_arr = arr_box.text_frame
            p_arr = tf_arr.paragraphs[0]
            p_arr.text = "➔"
            p_arr.font.size = Pt(16)
            p_arr.font.bold = True
            p_arr.font.color.rgb = ACCENT_INDIGO
            p_arr.alignment = PP_ALIGN.CENTER

    # Bottom Architectural Highlights Box
    arch_notes = add_card(slide6, Inches(0.8), Inches(5.7), Inches(11.733), Inches(1.1))
    tf_an = arch_notes.text_frame
    tf_an.margin_left = tf_an.margin_top = tf_an.margin_right = tf_an.margin_bottom = Inches(0.2)
    tf_an.word_wrap = True

    p_an1 = tf_an.paragraphs[0]
    p_an1.text = "🛡️ Key Architectural Guardrails:"
    p_an1.font.size = Pt(12)
    p_an1.font.bold = True
    p_an1.font.color.rgb = ACCENT_GREEN

    p_an2 = tf_an.add_paragraph()
    p_an2.text = "• 100% Offline execution via REST API (http://172.17.0.1:11434)  • Deterministic fallback parser to prevent model refusals  • Python-level ground-truth validation gates"
    p_an2.font.size = Pt(11)
    p_an2.font.color.rgb = TEXT_MUTED
    p_an2.space_before = Pt(4)

    # =========================================================================
    # SLIDE 7: Model, Tools & Jetson Deployment
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_header(slide7, "Model, Tools & Jetson Deployment Tech Stack")

    cols = [
        ("🤖 Model Used", [
            ("Model:", "Meta Llama 3.2 (1B Parameter)"),
            ("Type:", "Small Language Model (SLM)"),
            ("Embeddings:", "nomic-embed-text"),
            ("Quantization:", "Q4_K_M 4-bit Quantization"),
            ("Context Window:", "1024 Tokens (Optimized)"),
        ], ACCENT_INDIGO),
        ("🛠️ Tools & Stack", [
            ("Runtime:", "Python 3.14 & Ollama Local Server"),
            ("Framework:", "LangChain & Streamlit UI"),
            ("Vector Store:", "FAISS Vector Database"),
            ("Regex Engine:", "OCR Character Repair Engine"),
            ("Validation:", "Python Option Shuffler & Gate"),
        ], ACCENT_GREEN),
        ("⚡ Deployment on Jetson", [
            ("Hardware:", "NVIDIA Jetson Board (Orin / Nano)"),
            ("OS / SDK:", "JetPack / L4T CUDA Acceleration"),
            ("VRAM Footprint:", "~1.8 GB VRAM / RAM"),
            ("Optimization:", "Single-Pass Chunk Generation"),
            ("Network Mode:", "100% Air-Gapped Offline"),
        ], ACCENT_PURPLE),
    ]

    c_width = Inches(3.7)
    c_gap = Inches(0.3)

    for i, (title, items, color) in enumerate(cols):
        card = add_card(slide7, Inches(0.8) + i * (c_width + c_gap), Inches(1.6), c_width, Inches(5.2))
        tf = card.text_frame
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.25)
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = color

        for label, val in items:
            p = tf.add_paragraph()
            p.space_before = Pt(12)
            run1 = p.add_run()
            run1.text = f"• {label} "
            run1.font.bold = True
            run1.font.size = Pt(11)
            run1.font.color.rgb = TEXT_MUTED
            run2 = p.add_run()
            run2.text = val
            run2.font.bold = True
            run2.font.size = Pt(11)
            run2.font.color.rgb = TEXT_MAIN

    # =========================================================================
    # SLIDE 8: Live Demo Flow
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_header(slide8, "Live Demo Flow — Step-by-Step Execution")

    steps_demo = [
        ("Step 1: Document Ingestion", "Upload course PDF (e.g. Operating Systems Lecture Notes) into EduMind sidebar."),
        ("Step 2: Range & Topic Selection", "Select target page range (e.g. Pages 1–10) and optional custom topic focus."),
        ("Step 3: Pagewise Notes Trigger", "Click 'Pagewise Notes' — Receive instant, un-truncated markdown notes with headings."),
        ("Step 4: Practice Quiz Generation", "Click 'Practice Quiz' — System generates direct conceptual MCQs without fill-in gaps."),
        ("Step 5: Grading & Page Citation", "Submit quiz answers — Receive immediate grading with exact source page citations."),
        ("Step 6: Offline Jetson Verification", "Disconnect internet cable to demonstrate 100% Edge-AI operation on NVIDIA Jetson."),
    ]

    d_width = Inches(5.7)
    d_height = Inches(1.6)

    coords_d = [
        (Inches(0.8), Inches(1.6)),
        (Inches(6.833), Inches(1.6)),
        (Inches(0.8), Inches(3.4)),
        (Inches(6.833), Inches(3.4)),
        (Inches(0.8), Inches(5.2)),
        (Inches(6.833), Inches(5.2)),
    ]

    for i, (title, desc) in enumerate(steps_demo):
        x, y = coords_d[i]
        card = add_card(slide8, x, y, d_width, d_height)
        tf = card.text_frame
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.2)
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = ACCENT_GREEN

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_MAIN
        p_d.space_before = Pt(4)

    # =========================================================================
    # SLIDE 9: Key Metrics & Results
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9)
    add_header(slide9, "Key Metrics & Results — Benchmark Evaluation")

    metrics = [
        ("⚡ 1.2s - 2.5s", "Inference Latency", ACCENT_GREEN),
        ("🎯 92%", "Conceptual Accuracy", ACCENT_INDIGO),
        ("💾 ~1.8 GB", "VRAM / RAM Footprint", ACCENT_BLUE),
        ("🌐 100%", "Offline Independence", ACCENT_PURPLE),
        ("🧪 48 / 50", "PDF Benchmarks Passed", RGBColor(245, 158, 11)),
        ("🗣️ Multi-Lang", "Extensible Embeddings", RGBColor(236, 72, 153)),
    ]

    m_width = Inches(3.7)
    m_height = Inches(1.3)
    m_coords = [
        (Inches(0.8), Inches(1.5)),
        (Inches(4.8), Inches(1.5)),
        (Inches(8.8), Inches(1.5)),
        (Inches(0.8), Inches(3.0)),
        (Inches(4.8), Inches(3.0)),
        (Inches(8.8), Inches(3.0)),
    ]

    for i, (val, label, color) in enumerate(metrics):
        x, y = m_coords[i]
        card = add_card(slide9, x, y, m_width, m_height)
        tf = card.text_frame
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.15)
        tf.word_wrap = True

        p_v = tf.paragraphs[0]
        p_v.text = val
        p_v.font.size = Pt(20)
        p_v.font.bold = True
        p_v.font.color.rgb = color
        p_v.alignment = PP_ALIGN.CENTER

        p_l = tf.add_paragraph()
        p_l.text = label
        p_l.font.size = Pt(11)
        p_l.font.color.rgb = TEXT_MUTED
        p_l.alignment = PP_ALIGN.CENTER
        p_l.space_before = Pt(2)

    # Before vs After Table Card
    t_card = add_card(slide9, Inches(0.8), Inches(4.5), Inches(11.733), Inches(2.4))
    tf_t = t_card.text_frame
    tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = Inches(0.2)
    tf_t.word_wrap = True

    p_tt = tf_t.paragraphs[0]
    p_tt.text = "📊 Before vs. After EduMind Optimization:"
    p_tt.font.size = Pt(14)
    p_tt.font.bold = True
    p_tt.font.color.rgb = ACCENT_INDIGO

    # Table inside slide
    rows, cols = 3, 4
    left, top, width, height = Inches(1.0), Inches(5.0), Inches(11.333), Inches(1.7)
    table_shape = slide9.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(3.0)
    table.columns[2].width = Inches(3.0)
    table.columns[3].width = Inches(2.833)

    headers = ["Feature Metric", "Standard Cloud RAG", "EduMind Edge-AI Engine", "Improvement Impact"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30, 27, 75)
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ACCENT_INDIGO

    table_data = [
        ["Network Dependency", "Requires 100% Cloud Internet", "100% Local / Offline on Jetson", "Zero API Downtime"],
        ["MCQ Quality & OCR", "Verbatim fill-in-blanks & broken text", "Conceptual MCQs & Regex Repair", "92% Accuracy Pass Rate"],
    ]

    for i, row in enumerate(table_data):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_MAIN if j == 2 else TEXT_MUTED

    # =========================================================================
    # SLIDE 10: Novelty, Conclusion & Next Steps
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10)
    add_header(slide10, "Novelty, Conclusion & Future Roadmap")

    n_cols = [
        ("✨ What Makes It Novel?", "Solves the 1B SLM instruction-overload problem by delegating option shuffling and validation to Python while running regex text repair prior to vector inference.", ACCENT_INDIGO),
        ("🎓 Key Conclusion", "Demonstrates that lightweight 1B models deployed on edge hardware can outperform cloud APIs in reliability and domain-specific assessment tasks when paired with strict guardrails.", ACCENT_GREEN),
        ("🚀 Future Roadmap", "1. Voice-guided quiz audio narration.\n2. Multimodal lecture diagram parsing.\n3. Quantized 3B parameter model support on Jetson Orin.", ACCENT_PURPLE),
    ]

    n_width = Inches(3.7)
    n_gap = Inches(0.3)

    for i, (title, desc, color) in enumerate(n_cols):
        card = add_card(slide10, Inches(0.8) + i * (n_width + n_gap), Inches(1.6), n_width, Inches(4.0))
        tf = card.text_frame
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.25)
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = color

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = TEXT_MAIN
        p_d.space_before = Pt(12)

    # Closing Q&A Box
    close_card = add_card(slide10, Inches(0.8), Inches(5.8), Inches(11.733), Inches(1.0), bg_color=RGBColor(30, 27, 75), border_color=ACCENT_GREEN)
    tf_cl = close_card.text_frame
    tf_cl.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_cl = tf_cl.paragraphs[0]
    p_cl.text = "💬 Thank You! Happy to take questions, live feedback, and evaluation comments."
    p_cl.font.size = Pt(14)
    p_cl.font.bold = True
    p_cl.font.color.rgb = ACCENT_GREEN
    p_cl.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 11: Project Evaluation Rubric
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11)
    add_header(slide11, "Project Evaluation Rubric — EDGEMINDS 2026")

    # Table Card
    r_card = add_card(slide11, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.4))
    tf_r = r_card.text_frame
    tf_r.margin_left = tf_r.margin_top = tf_r.margin_right = tf_r.margin_bottom = Inches(0.2)

    p_rt = tf_r.paragraphs[0]
    p_rt.text = "📋 EDGEMINDS Internship 2026 Evaluation Reference Matrix:"
    p_rt.font.size = Pt(14)
    p_rt.font.bold = True
    p_rt.font.color.rgb = ACCENT_INDIGO

    rows, cols = 6, 4
    left, top, width, height = Inches(1.0), Inches(2.1), Inches(11.333), Inches(4.5)
    t_shape = slide11.shapes.add_table(rows, cols, left, top, width, height)
    t_rubric = t_shape.table

    t_rubric.columns[0].width = Inches(2.2)
    t_rubric.columns[1].width = Inches(1.3)
    t_rubric.columns[2].width = Inches(5.833)
    t_rubric.columns[3].width = Inches(2.0)

    r_headers = ["Evaluation Criteria", "Weightage", "Description & Implementation Standard", "EduMind Compliance"]
    for j, h in enumerate(r_headers):
        cell = t_rubric.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30, 27, 75)
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ACCENT_INDIGO

    r_data = [
        ["Functionality", "30 Wt.", "Works reliably on Jetson board, completes core task without crashing.", "100% Functional on Jetson"],
        ["Innovation", "25 Wt.", "Novel Edge-RAG architecture beyond basic templates, featuring regex repair.", "Novel Dual Quality Pipeline"],
        ["Edge-Readiness", "20 Wt.", "100% offline execution, optimized VRAM/memory footprint (<1.8 GB).", "Air-Gapped CUDA L4T"],
        ["Technical Quality", "15 Wt.", "Clean Python 3.14 code, regex repair engines, and ground-truth validation.", "Production Code Structure"],
        ["Presentation & Demo", "10 Wt.", "Clear story, live Jetson demonstration, and thorough documentation.", "Live Demo & Backup Log"],
    ]

    for i, row in enumerate(r_data):
        for j, val in enumerate(row):
            cell = t_rubric.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(10)
            p.font.color.rgb = ACCENT_GREEN if j == 3 else (TEXT_MAIN if j == 0 else TEXT_MUTED)
            if j == 1:
                p.font.bold = True

    # Save presentation
    import sys
    sys.stdout.reconfigure(encoding='utf-8')
    output_path = "edumind_final_presentation.pptx"
    prs.save(output_path)
    print(f"Successfully generated '{output_path}' with 11 widescreen slides.")


if __name__ == "__main__":
    create_edumind_deck()
