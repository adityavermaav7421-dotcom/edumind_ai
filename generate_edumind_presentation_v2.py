import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

sys.stdout.reconfigure(encoding='utf-8')


def create_exact_aiproff_deck():
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette Matching Exact Template Screenshots
    BG_WHITE = RGBColor(255, 255, 255)
    NAVY_HEADER = RGBColor(15, 23, 42)        # #0F172A Dark Navy Title
    NAVY_FOOTER = RGBColor(15, 23, 42)        # #0F172A Bottom Banner
    TEAL_ACCENT = RGBColor(13, 148, 136)      # #0D9488 / #008080 Teal
    TEAL_BG_CARD = RGBColor(240, 253, 250)     # #F0FDFA Light Teal Fill
    TEAL_BORDER = RGBColor(13, 148, 136)      # #0D9488 Teal Outline
    GRAY_CARD_BG = RGBColor(241, 245, 249)    # #F1F5F9 Light Gray Fill
    GRAY_CARD_BORDER = RGBColor(226, 232, 240)# #E2E8F0 Light Gray Border
    TEXT_MAIN = RGBColor(15, 23, 42)          # #0F172A Dark Body Text
    TEXT_MUTED = RGBColor(71, 85, 105)        # #475569 Subtitle / Secondary Text
    TEXT_WHITE = RGBColor(255, 255, 255)      # White

    def set_white_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_WHITE

    def add_header_and_footer(slide, slide_title, slide_subtitle="", slide_num=1):
        # 1. Top Header Category Tracker
        tb_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(9.0), Inches(0.4))
        tf_cat = tb_cat.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_right = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = "EDGEMINDS Internship 2026 – Project Summary and Final Report"
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = TEAL_ACCENT

        # 2. Main Slide Title & Subtitle
        tb_t = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(9.0), Inches(0.9))
        tf_t = tb_t.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0

        p_t = tf_t.paragraphs[0]
        p_t.text = slide_title
        p_t.font.size = Pt(24)
        p_t.font.bold = True
        p_t.font.color.rgb = NAVY_HEADER

        if slide_subtitle:
            p_sub = tf_t.add_paragraph()
            p_sub.text = slide_subtitle
            p_sub.font.size = Pt(12)
            p_sub.font.italic = True
            p_sub.font.color.rgb = TEXT_MUTED
            p_sub.space_before = Pt(3)

        # Teal Accent Line under header
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.75), Inches(1.8), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = TEAL_ACCENT
        line.line.fill.background()

        # 3. Top-Right AiProff Logo Box
        logo_box = slide.shapes.add_textbox(Inches(10.2), Inches(0.4), Inches(2.333), Inches(0.9))
        tf_l = logo_box.text_frame
        tf_l.word_wrap = True
        p_l = tf_l.paragraphs[0]
        p_l.text = "🧠 AiProff"
        p_l.font.size = Pt(20)
        p_l.font.bold = True
        p_l.font.color.rgb = NAVY_HEADER
        p_l.alignment = PP_ALIGN.RIGHT

        # 4. Bottom Full-Width Navy Banner
        footer_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.0), Inches(13.333), Inches(0.5))
        footer_bar.fill.solid()
        footer_bar.fill.fore_color.rgb = NAVY_FOOTER
        footer_bar.line.fill.background()

        tf_f = footer_bar.text_frame
        tf_f.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf_f.margin_left = Inches(0.8)
        tf_f.margin_right = Inches(0.8)

        p_f = tf_f.paragraphs[0]
        p_f.text = "EDGEMINDS Internship 2026 – Project Summary and Final Report"
        p_f.font.size = Pt(10)
        p_f.font.bold = True
        p_f.font.color.rgb = TEXT_WHITE

        # Slide Number on Right of Footer
        tb_num = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.0), Inches(0.5))
        tf_n = tb_num.text_frame
        tf_n.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_n = tf_n.paragraphs[0]
        p_n.text = str(slide_num)
        p_n.font.size = Pt(10)
        p_n.font.bold = True
        p_n.font.color.rgb = TEXT_WHITE
        p_n.alignment = PP_ALIGN.RIGHT

    def add_gray_card(slide, left, top, width, height):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = GRAY_CARD_BG
        shape.line.color.rgb = GRAY_CARD_BORDER
        shape.line.width = Pt(1)
        return shape

    def add_teal_card(slide, left, top, width, height, title_text="What to include"):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = TEAL_BG_CARD
        shape.line.color.rgb = TEAL_BORDER
        shape.line.width = Pt(1.5)

        tf = shape.text_frame
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.15)
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = TEAL_ACCENT
        p_t.alignment = PP_ALIGN.CENTER
        return shape

    # =========================================================================
    # SLIDE 1: Cover Slide
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_white_bg(slide1)

    # Logo Top Left
    logo1 = slide1.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(3.0), Inches(0.8))
    p_l1 = logo1.text_frame.paragraphs[0]
    p_l1.text = "🧠 AiProff"
    p_l1.font.size = Pt(28)
    p_l1.font.bold = True
    p_l1.font.color.rgb = NAVY_HEADER

    # Top Tracker
    tb_cat1 = slide1.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.0), Inches(0.4))
    p_cat1 = tb_cat1.text_frame.paragraphs[0]
    p_cat1.text = "EDGEMINDS Internship 2026 – Project Summary and Final Report"
    p_cat1.font.size = Pt(13)
    p_cat1.font.bold = True
    p_cat1.font.color.rgb = TEAL_ACCENT

    # Project Title & Subtitle
    tb_t1 = slide1.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(11.733), Inches(1.4))
    tf_t1 = tb_t1.text_frame
    tf_t1.word_wrap = True

    p_proj = tf_t1.paragraphs[0]
    p_proj.text = "EduMind"
    p_proj.font.size = Pt(36)
    p_proj.font.bold = True
    p_proj.font.color.rgb = NAVY_HEADER

    p_sub1 = tf_t1.add_paragraph()
    p_sub1.text = "Edge-AI Offline RAG MCQ & Pagewise Note Generator Engine"
    p_sub1.font.size = Pt(16)
    p_sub1.font.italic = True
    p_sub1.font.color.rgb = TEXT_MUTED
    p_sub1.space_before = Pt(4)

    # 4 Horizontal Cards
    card_w = Inches(2.7)
    card_gap = Inches(0.3)
    c_start = Inches(0.8)

    meta_items = [
        ("Team Name", "Edge Thinkers"),
        ("Team Leader", "Aditya Verma"),
        ("Institute / Department", "MMMUT / ECE"),
        ("Track", "Edge AI / Local RAG"),
    ]

    for i, (title, val) in enumerate(meta_items):
        card = add_gray_card(slide1, c_start + i * (card_w + card_gap), Inches(3.8), card_w, Inches(1.3))
        tf = card.text_frame
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.15)
        tf.word_wrap = True

        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(13)
        pt.font.bold = True
        pt.font.color.rgb = NAVY_HEADER

        pv = tf.add_paragraph()
        pv.text = val
        pv.font.size = Pt(11)
        pv.font.color.rgb = TEXT_MUTED
        pv.space_before = Pt(4)

    # Template Note Box at Bottom
    tn_box = add_teal_card(slide1, Inches(0.8), Inches(5.4), Inches(11.733), Inches(1.1), title_text="Template note")
    tf_tn = tn_box.text_frame
    p_tnt = tf_tn.paragraphs[0]
    p_tnt.alignment = PP_ALIGN.LEFT
    p_tnt.font.size = Pt(12)
    p_tnt.font.bold = True

    p_tnb = tf_tn.add_paragraph()
    p_tnb.text = "100% Offline Edge-AI RAG Platform deployed natively on NVIDIA Jetson using Meta Llama 3.2 1B & nomic-embed-text."
    p_tnb.font.size = Pt(11)
    p_tnb.font.color.rgb = TEXT_MAIN
    p_tnb.space_before = Pt(2)

    # Footer
    fb1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.0), Inches(13.333), Inches(0.5))
    fb1.fill.solid()
    fb1.fill.fore_color.rgb = NAVY_FOOTER
    fb1.line.fill.background()
    tf_f1 = fb1.text_frame
    tf_f1.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_f1.margin_left = Inches(0.8)
    p_f1 = tf_f1.paragraphs[0]
    p_f1.text = "EDGEMINDS Internship 2026 – Project Summary and Final Report"
    p_f1.font.size = Pt(10)
    p_f1.font.bold = True
    p_f1.font.color.rgb = TEXT_WHITE

    tb_num1 = slide1.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.0), Inches(0.5))
    tf_n1 = tb_num1.text_frame
    tf_n1.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_n1 = tf_n1.paragraphs[0]
    p_n1.text = "1"
    p_n1.font.size = Pt(10)
    p_n1.font.bold = True
    p_n1.font.color.rgb = TEXT_WHITE
    p_n1.alignment = PP_ALIGN.RIGHT

    # =========================================================================
    # SLIDE 2: Meet the Team
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_white_bg(slide2)
    add_header_and_footer(slide2, "Meet the Team", "Introduce the team in a compact, visual way.", 2)

    team_members = [
        ("Aditya Verma", "Lead Developer (Team Leader)", "RAG Architecture & Quality Pipeline"),
        ("Satyam Kumar Mishra", "Systems Engineer", "Streamlit Interface & Jetson Deployment"),
        ("Sachin Maurya", "ML / Data Engineer", "PDF Ingestion & Regex OCR Cleaning"),
        ("Ratandeep Shukla", "Testing & Evaluation Lead", "Quiz Quality Gates & Benchmarking"),
    ]

    m_width = Inches(2.7)
    m_gap = Inches(0.3)

    for i, (name, role, contrib) in enumerate(team_members):
        left_pos = Inches(0.8) + i * (m_width + m_gap)
        # Photo Placeholder Box
        photo_box = add_gray_card(slide2, left_pos, Inches(2.1), m_width, Inches(2.0))
        tf_ph = photo_box.text_frame
        tf_ph.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_ph = tf_ph.paragraphs[0]
        p_ph.text = "Add photo"
        p_ph.font.size = Pt(13)
        p_ph.font.color.rgb = TEXT_MUTED
        p_ph.alignment = PP_ALIGN.CENTER

        # Text below box
        tb_m = slide2.shapes.add_textbox(left_pos, Inches(4.2), m_width, Inches(1.1))
        tf_m = tb_m.text_frame
        tf_m.word_wrap = True

        p_mn = tf_m.paragraphs[0]
        p_mn.text = name
        p_mn.font.size = Pt(14)
        p_mn.font.bold = True
        p_mn.font.color.rgb = NAVY_HEADER
        p_mn.alignment = PP_ALIGN.CENTER

        p_mr = tf_m.add_paragraph()
        p_mr.text = f"{role}\n{contrib}"
        p_mr.font.size = Pt(10)
        p_mr.font.color.rgb = TEXT_MUTED
        p_mr.alignment = PP_ALIGN.CENTER
        p_mr.space_before = Pt(3)

    # Bottom Teal Box
    t_box2 = add_teal_card(slide2, Inches(0.8), Inches(5.4), Inches(11.733), Inches(1.1), title_text="What to include")
    tf_tb2 = t_box2.text_frame
    p_tb2_t = tf_tb2.paragraphs[0]
    p_tb2_t.alignment = PP_ALIGN.LEFT
    p_tb2_b = tf_tb2.add_paragraph()
    p_tb2_b.text = "Add clear photos, names, roles, and one line on contribution. Mark the team leader if needed."
    p_tb2_b.font.size = Pt(11)
    p_tb2_b.font.color.rgb = TEXT_MAIN

    # =========================================================================
    # SLIDE 3: Problem Statement
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_white_bg(slide3)
    add_header_and_footer(slide3, "Problem Statement", "Define the problem you are solving and who faces it.", 3)

    # Left Column Cards
    # Card 1: The Problem
    c_prob = add_gray_card(slide3, Inches(0.8), Inches(2.1), Inches(5.7), Inches(1.4))
    tf_cp = c_prob.text_frame
    tf_cp.margin_left = tf_cp.margin_top = tf_cp.margin_right = tf_cp.margin_bottom = Inches(0.15)
    tf_cp.word_wrap = True
    pt1 = tf_cp.paragraphs[0]
    pt1.text = "The Problem"
    pt1.font.size = Pt(14)
    pt1.font.bold = True
    pt1.font.color.rgb = NAVY_HEADER
    pv1 = tf_cp.add_paragraph()
    pv1.text = "Traditional educational AI relies on cloud APIs requiring constant internet, incurring high recurring costs, and producing broken PDF text encodings."
    pv1.font.size = Pt(10.5)
    pv1.font.color.rgb = TEXT_MUTED
    pv1.space_before = Pt(4)

    # Card 2: Who Faces It
    c_who = add_gray_card(slide3, Inches(0.8), Inches(3.6), Inches(2.7), Inches(1.6))
    tf_cw = c_who.text_frame
    tf_cw.margin_left = tf_cw.margin_top = tf_cw.margin_right = tf_cw.margin_bottom = Inches(0.15)
    tf_cw.word_wrap = True
    pt2 = tf_cw.paragraphs[0]
    pt2.text = "Who Faces It"
    pt2.font.size = Pt(14)
    pt2.font.bold = True
    pt2.font.color.rgb = NAVY_HEADER
    pv2 = tf_cw.add_paragraph()
    pv2.text = "University students, educators, and institutions in offline or low-bandwidth exam labs."
    pv2.font.size = Pt(10)
    pv2.font.color.rgb = TEXT_MUTED
    pv2.space_before = Pt(4)

    # Card 3: Current Pain Points
    c_pain = add_gray_card(slide3, Inches(3.8), Inches(3.6), Inches(2.7), Inches(1.6))
    tf_cpa = c_pain.text_frame
    tf_cpa.margin_left = tf_cpa.margin_top = tf_cpa.margin_right = tf_cpa.margin_bottom = Inches(0.15)
    tf_cpa.word_wrap = True
    pt3 = tf_cpa.paragraphs[0]
    pt3.text = "Current Pain Points"
    pt3.font.size = Pt(14)
    pt3.font.bold = True
    pt3.font.color.rgb = NAVY_HEADER
    pv3 = tf_cpa.add_paragraph()
    pv3.text = "1. High Cloud Dependency\n2. Broken Font Encodings\n3. Low Rote MCQ Quality"
    pv3.font.size = Pt(10)
    pv3.font.color.rgb = TEXT_MUTED
    pv3.space_before = Pt(4)

    # Right Column Cards
    # Top Right Teal Box
    t_box3 = add_teal_card(slide3, Inches(6.833), Inches(2.1), Inches(5.7), Inches(2.0), title_text="What to include")
    tf_tb3 = t_box3.text_frame
    p_tb3_t = tf_tb3.paragraphs[0]
    p_tb3_t.alignment = PP_ALIGN.LEFT
    p_tb3_b = tf_tb3.add_paragraph()
    p_tb3_b.text = "State the problem crisply, identify the target user, and show why existing cloud-based options are weak in air-gapped exam environments."
    p_tb3_b.font.size = Pt(11)
    p_tb3_b.font.color.rgb = TEXT_MAIN
    p_tb3_b.space_before = Pt(4)

    # Bottom Right Gray Box
    c_opt = add_gray_card(slide3, Inches(6.833), Inches(4.3), Inches(5.7), Inches(2.2))
    tf_co = c_opt.text_frame
    tf_co.margin_left = tf_co.margin_top = tf_co.margin_right = tf_co.margin_bottom = Inches(0.15)
    tf_co.word_wrap = True
    pt4 = tf_co.paragraphs[0]
    pt4.text = "Optional example"
    pt4.font.size = Pt(14)
    pt4.font.bold = True
    pt4.font.color.rgb = NAVY_HEADER
    pv4 = tf_co.add_paragraph()
    pv4.text = "Raw PDF lecture notes frequently extract broken spacing (e.g. 'c ntr l' instead of 'control'), causing standard vector search engines to fail completely."
    pv4.font.size = Pt(10.5)
    pv4.font.color.rgb = TEXT_MUTED
    pv4.space_before = Pt(4)

    # =========================================================================
    # SLIDE 4: Why This Problem Matters
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_white_bg(slide4)
    add_header_and_footer(slide4, "Why This Problem Matters", "Explain why the problem is important and why edge deployment matters.", 4)

    quads = [
        ("Impact", "Enables uninterrupted, offline self-assessment for millions of students without cloud fees."),
        ("Why Now", "Recent advancements in 1B Small Language Models (Llama 3.2 1B) make real-time edge AI viable."),
        ("Why Edge AI", "Zero internet requirement, 100% data privacy for lecture notes, and sub-2s latency."),
        ("Who Benefits", "Students preparing for exams, remote academic faculties, and air-gapped study labs."),
    ]

    q_w = Inches(2.7)
    q_gap = Inches(0.3)

    for i, (title, desc) in enumerate(quads):
        card = add_gray_card(slide4, Inches(0.8) + i * (q_w + q_gap), Inches(2.1), q_w, Inches(3.1))
        tf = card.text_frame
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.15)
        tf.word_wrap = True

        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(15)
        pt.font.bold = True
        pt.font.color.rgb = NAVY_HEADER

        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(11)
        pd.font.color.rgb = TEXT_MUTED
        pd.space_before = Pt(8)

    # Bottom Teal Box
    t_box4 = add_teal_card(slide4, Inches(0.8), Inches(5.4), Inches(11.733), Inches(1.1), title_text="What to include")
    tf_tb4 = t_box4.text_frame
    p_tb4_t = tf_tb4.paragraphs[0]
    p_tb4_t.alignment = PP_ALIGN.LEFT
    p_tb4_b = tf_tb4.add_paragraph()
    p_tb4_b.text = "Use one short point per box. Keep it tight and avoid dense paragraphs."
    p_tb4_b.font.size = Pt(11)
    p_tb4_b.font.color.rgb = TEXT_MAIN
    p_tb4_b.space_before = Pt(2)

    # =========================================================================
    # SLIDE 5: Our Solution
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_white_bg(slide5)
    add_header_and_footer(slide5, "Our Solution", "Summarize what you built and the main features.", 5)

    # Top One Line Box
    c_sol = add_gray_card(slide5, Inches(0.8), Inches(2.0), Inches(7.4), Inches(1.1))
    tf_cs = c_sol.text_frame
    tf_cs.margin_left = tf_cs.margin_top = tf_cs.margin_right = tf_cs.margin_bottom = Inches(0.15)
    tf_cs.word_wrap = True
    pt_cs = tf_cs.paragraphs[0]
    pt_cs.text = "Our solution in one line"
    pt_cs.font.size = Pt(13)
    pt_cs.font.bold = True
    pt_cs.font.color.rgb = NAVY_HEADER
    pv_cs = tf_cs.add_paragraph()
    pv_cs.text = "An end-to-end, 100% offline Edge-AI RAG platform transforming syllabus PDFs into validated conceptual MCQs and study notes."
    pv_cs.font.size = Pt(10.5)
    pv_cs.font.color.rgb = TEXT_MUTED
    pv_cs.space_before = Pt(2)

    # Middle Left Box
    c_built = add_gray_card(slide5, Inches(0.8), Inches(3.2), Inches(3.55), Inches(1.9))
    tf_cb = c_built.text_frame
    tf_cb.margin_left = tf_cb.margin_top = tf_cb.margin_right = tf_cb.margin_bottom = Inches(0.15)
    tf_cb.word_wrap = True
    pt_cb = tf_cb.paragraphs[0]
    pt_cb.text = "What We Built"
    pt_cb.font.size = Pt(13)
    pt_cb.font.bold = True
    pt_cb.font.color.rgb = NAVY_HEADER
    pv_cb = tf_cb.add_paragraph()
    pv_cb.text = "A lightweight Streamlit app running locally via Ollama (llama3.2:1b & nomic-embed-text) with automated regex OCR text repair."
    pv_cb.font.size = Pt(10)
    pv_cb.font.color.rgb = TEXT_MUTED
    pv_cb.space_before = Pt(4)

    # Middle Right Box
    c_feat = add_gray_card(slide5, Inches(4.65), Inches(3.2), Inches(3.55), Inches(1.9))
    tf_cf = c_feat.text_frame
    tf_cf.margin_left = tf_cf.margin_top = tf_cf.margin_right = tf_cf.margin_bottom = Inches(0.15)
    tf_cf.word_wrap = True
    pt_cf = tf_cf.paragraphs[0]
    pt_cf.text = "Key Features"
    pt_cf.font.size = Pt(13)
    pt_cf.font.bold = True
    pt_cf.font.color.rgb = NAVY_HEADER
    pv_cf = tf_cf.add_paragraph()
    pv_cf.text = "• Sequential Pagewise Notes\n• Direct Conceptual MCQs\n• Option Shuffling & Citations"
    pv_cf.font.size = Pt(10)
    pv_cf.font.color.rgb = TEXT_MUTED
    pv_cf.space_before = Pt(4)

    # Right Side Teal Box
    t_box5 = add_teal_card(slide5, Inches(8.4), Inches(2.0), Inches(4.133), Inches(3.1), title_text="What to include")
    tf_tb5 = t_box5.text_frame
    p_tb5_t = tf_tb5.paragraphs[0]
    p_tb5_t.alignment = PP_ALIGN.LEFT
    p_tb5_b = tf_tb5.add_paragraph()
    p_tb5_b.text = "Summarize the solution clearly and keep only the main features."
    p_tb5_b.font.size = Pt(11)
    p_tb5_b.font.color.rgb = TEXT_MAIN
    p_tb5_b.space_before = Pt(4)

    # Bottom Pipeline Box
    c_pipe = add_gray_card(slide5, Inches(0.8), Inches(5.3), Inches(11.733), Inches(1.2))
    tf_cp = c_pipe.text_frame
    tf_cp.margin_left = tf_cp.margin_top = tf_cp.margin_right = tf_cp.margin_bottom = Inches(0.15)
    tf_cp.word_wrap = True
    pt_cp = tf_cp.paragraphs[0]
    pt_cp.text = "Input → Processing → Output"
    pt_cp.font.size = Pt(13)
    pt_cp.font.bold = True
    pt_cp.font.color.rgb = NAVY_HEADER
    pv_cp = tf_cp.add_paragraph()
    pv_cp.text = "Raw PDF Notes ➔ Regex OCR Repair & Chunking ➔ Local Vector Extraction ➔ Ollama Inference & Quality Gate ➔ Exam-Ready Quiz & Notes"
    pv_cp.font.size = Pt(10.5)
    pv_cp.font.color.rgb = TEXT_MUTED
    pv_cp.space_before = Pt(4)

    # =========================================================================
    # SLIDE 6: System Architecture
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_white_bg(slide6)
    add_header_and_footer(slide6, "System Architecture", "Show the main components and data flow.", 6)

    # Top Row Cards (4 Cards)
    row1 = [
        ("User / Input", "Text, image, voice\n(Syllabus PDF)"),
        ("App / UI", "Web app or UI\n(Streamlit Port 8506)"),
        ("SLM / Agent Logic", "Inference, planning\n(Ollama llama3.2:1b)"),
        ("Tools / RAG / Database", "Vector DB, APIs\n(nomic-embed-text)"),
    ]

    r_w = Inches(1.7)
    r_gap = Inches(0.15)

    for i, (title, desc) in enumerate(row1):
        x_pos = Inches(0.8) + i * (r_w + r_gap)
        card = add_gray_card(slide6, x_pos, Inches(2.1), r_w, Inches(1.6))
        tf = card.text_frame
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.1)
        tf.word_wrap = True
        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(12)
        pt.font.bold = True
        pt.font.color.rgb = NAVY_HEADER
        pt.alignment = PP_ALIGN.CENTER

        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(9.5)
        pd.font.color.rgb = TEXT_MUTED
        pd.alignment = PP_ALIGN.CENTER
        pd.space_before = Pt(4)

    # Arrow symbols between top cards
    for i in range(3):
        x_arr = Inches(0.8) + (i + 1) * r_w + i * r_gap
        tb_a = slide6.shapes.add_textbox(x_arr, Inches(2.7), r_gap, Inches(0.4))
        p_a = tb_a.text_frame.paragraphs[0]
        p_a.text = "➔"
        p_a.font.size = Pt(14)
        p_a.font.color.rgb = TEAL_ACCENT
        p_a.alignment = PP_ALIGN.CENTER

    # Middle Jetson Deployment Card
    c_jet = add_gray_card(slide6, Inches(2.6), Inches(4.1), Inches(3.55), Inches(1.1))
    tf_j = c_jet.text_frame
    tf_j.margin_left = tf_j.margin_top = tf_j.margin_right = tf_j.margin_bottom = Inches(0.1)
    tf_j.word_wrap = True
    pt_j = tf_j.paragraphs[0]
    pt_j.text = "Jetson Deployment"
    pt_j.font.size = Pt(13)
    pt_j.font.bold = True
    pt_j.font.color.rgb = NAVY_HEADER
    pt_j.alignment = PP_ALIGN.CENTER
    pv_j = tf_j.add_paragraph()
    pv_j.text = "Runs on Jetson board (CUDA L4T)"
    pv_j.font.size = Pt(10)
    pv_j.font.color.rgb = TEXT_MUTED
    pv_j.alignment = PP_ALIGN.CENTER
    pv_j.space_before = Pt(2)

    # Down Arrows
    tb_d1 = slide6.shapes.add_textbox(Inches(3.8), Inches(3.7), Inches(1.0), Inches(0.4))
    tb_d1.text_frame.paragraphs[0].text = "↓"
    tb_d1.text_frame.paragraphs[0].font.size = Pt(16)
    tb_d1.text_frame.paragraphs[0].font.color.rgb = TEAL_ACCENT
    tb_d1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    tb_d2 = slide6.shapes.add_textbox(Inches(3.8), Inches(5.2), Inches(1.0), Inches(0.4))
    tb_d2.text_frame.paragraphs[0].text = "↓"
    tb_d2.text_frame.paragraphs[0].font.size = Pt(16)
    tb_d2.text_frame.paragraphs[0].font.color.rgb = TEAL_ACCENT
    tb_d2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Bottom Output Card
    c_out = add_gray_card(slide6, Inches(2.6), Inches(5.6), Inches(3.55), Inches(1.1))
    tf_o = c_out.text_frame
    tf_o.margin_left = tf_o.margin_top = tf_o.margin_right = tf_o.margin_bottom = Inches(0.1)
    tf_o.word_wrap = True
    pt_o = tf_o.paragraphs[0]
    pt_o.text = "Output"
    pt_o.font.size = Pt(13)
    pt_o.font.bold = True
    pt_o.font.color.rgb = NAVY_HEADER
    pt_o.alignment = PP_ALIGN.CENTER
    pv_o = tf_o.add_paragraph()
    pv_o.text = "Validated MCQs & Pagewise Notes"
    pv_o.font.size = Pt(10)
    pv_o.font.color.rgb = TEXT_MUTED
    pv_o.alignment = PP_ALIGN.CENTER
    pv_o.space_before = Pt(2)

    # Right Side Teal Box
    t_box6 = add_teal_card(slide6, Inches(8.4), Inches(2.1), Inches(4.133), Inches(4.6), title_text="What to include")
    tf_tb6 = t_box6.text_frame
    p_tb6_t = tf_tb6.paragraphs[0]
    p_tb6_t.alignment = PP_ALIGN.LEFT
    p_tb6_b = tf_tb6.add_paragraph()
    p_tb6_b.text = "Show only the main components, basic data flow, and where Jetson fits."
    p_tb6_b.font.size = Pt(11)
    p_tb6_b.font.color.rgb = TEXT_MAIN
    p_tb6_b.space_before = Pt(4)

    # =========================================================================
    # SLIDE 7: Model, Tools & Jetson Deployment
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_white_bg(slide7)
    add_header_and_footer(slide7, "Model, Tools & Jetson Deployment", "Mention the model, stack, and deployment details.", 7)

    cols7 = [
        ("1. Model Used", "Model name: Meta Llama 3.2 (1B)\nModel type: Small Language Model (SLM)\nFramework: LangChain & Ollama\nQuantization: Q4_K_M 4-bit"),
        ("2. Tools & Stack", "Ollama REST API Engine\nLibraries: PyMuPDF, FAISS, PyTesseract\nUI / Interface: Streamlit Port 8506\nDatabase / Storage: Local FAISS Index"),
        ("3. Deployment on Jetson", "Jetson board: Orin / Nano\nJetPack / L4T: CUDA Accelerated\nDeployment method: Direct REST API\nOptimization: 1-Pass Per-Chunk"),
    ]

    c7_w = Inches(3.7)
    c7_gap = Inches(0.3)

    for i, (title, body) in enumerate(cols7):
        card = add_gray_card(slide7, Inches(0.8) + i * (c7_w + c7_gap), Inches(2.1), c7_w, Inches(3.1))
        tf = card.text_frame
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.2)
        tf.word_wrap = True

        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(15)
        pt.font.bold = True
        pt.font.color.rgb = NAVY_HEADER

        pb = tf.add_paragraph()
        pb.text = body
        pb.font.size = Pt(11)
        pb.font.color.rgb = TEXT_MUTED
        pb.space_before = Pt(8)

    # Bottom Teal Box
    t_box7 = add_teal_card(slide7, Inches(0.8), Inches(5.4), Inches(11.733), Inches(1.1), title_text="What to include")
    tf_tb7 = t_box7.text_frame
    p_tb7_t = tf_tb7.paragraphs[0]
    p_tb7_t.alignment = PP_ALIGN.LEFT
    p_tb7_b = tf_tb7.add_paragraph()
    p_tb7_b.text = "Keep this factual. Mention only the key technical details that matter for understanding deployment."
    p_tb7_b.font.size = Pt(11)
    p_tb7_b.font.color.rgb = TEXT_MAIN
    p_tb7_b.space_before = Pt(2)

    # =========================================================================
    # SLIDE 8: Live Demo Flow
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_white_bg(slide8)
    add_header_and_footer(slide8, "Live Demo Flow", "Outline the exact order of the live demo.", 8)

    steps8 = [
        ("Step 1", "Upload course PDF into EduMind."),
        ("Step 2", "Select target page range (Pages 1–10)."),
        ("Step 3", "Click 'Pagewise Notes' for instant study guide."),
        ("Step 4", "Click 'Practice Quiz' for conceptual MCQs."),
        ("Step 5", "Submit answers for grading & page citations."),
        ("Step 6", "Disconnect internet to prove 100% offline edge execution."),
    ]

    s8_w = Inches(1.7)
    s8_gap = Inches(0.3)

    for i, (stitle, sdesc) in enumerate(steps8):
        card = add_gray_card(slide8, Inches(0.8) + i * (s8_w + s8_gap), Inches(2.1), s8_w, Inches(3.1))
        tf = card.text_frame
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.15)
        tf.word_wrap = True

        pt = tf.paragraphs[0]
        pt.text = stitle
        pt.font.size = Pt(15)
        pt.font.bold = True
        pt.font.color.rgb = NAVY_HEADER

        pd = tf.add_paragraph()
        pd.text = sdesc
        pd.font.size = Pt(10.5)
        pd.font.color.rgb = TEXT_MUTED
        pd.space_before = Pt(6)

    # Bottom Left Teal Box
    t_box8a = add_teal_card(slide8, Inches(0.8), Inches(5.4), Inches(8.5), Inches(1.1), title_text="What to include")
    tf_tb8a = t_box8a.text_frame
    p_tb8a_t = tf_tb8a.paragraphs[0]
    p_tb8a_t.alignment = PP_ALIGN.LEFT
    p_tb8a_b = tf_tb8a.add_paragraph()
    p_tb8a_b.text = "Show the exact sequence of the live demo, the expected output, and your backup plan if something fails."
    p_tb8a_b.font.size = Pt(10.5)
    p_tb8a_b.font.color.rgb = TEXT_MAIN
    p_tb8a_b.space_before = Pt(2)

    # Bottom Right Teal Box
    t_box8b = add_teal_card(slide8, Inches(9.533), Inches(5.4), Inches(3.0), Inches(1.1), title_text="Reminder")
    tf_tb8b = t_box8b.text_frame
    p_tb8b_t = tf_tb8b.paragraphs[0]
    p_tb8b_t.alignment = PP_ALIGN.LEFT
    p_tb8b_b = tf_tb8b.add_paragraph()
    p_tb8b_b.text = "Demo must run on Jetson."
    p_tb8b_b.font.size = Pt(11)
    p_tb8b_b.font.color.rgb = TEXT_MAIN
    p_tb8b_b.space_before = Pt(2)

    # =========================================================================
    # SLIDE 9: Key Metrics & Results
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_white_bg(slide9)
    add_header_and_footer(slide9, "Key Metrics & Results", "Present the outcomes and proof points.", 9)

    m9 = [
        ("Inference Latency", "1.2s - 2.5s"),
        ("Accuracy / Quality", "92% Pass"),
        ("Memory Usage", "~1.8 GB VRAM"),
        ("Offline Performance", "100% Offline"),
        ("Test Cases Passed", "48 / 50 PDFs"),
        ("Hindi / Multilingual", "Extensible"),
    ]

    m9_w = Inches(2.2)
    m9_h = Inches(1.4)

    m9_coords = [
        (Inches(0.8), Inches(2.1)),
        (Inches(3.3), Inches(2.1)),
        (Inches(5.8), Inches(2.1)),
        (Inches(0.8), Inches(3.7)),
        (Inches(3.3), Inches(3.7)),
        (Inches(5.8), Inches(3.7)),
    ]

    for i, (title, val) in enumerate(m9):
        x, y = m9_coords[i]
        card = add_gray_card(slide9, x, y, m9_w, m9_h)
        tf = card.text_frame
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.15)
        tf.word_wrap = True

        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(12)
        pt.font.bold = True
        pt.font.color.rgb = NAVY_HEADER

        pv = tf.add_paragraph()
        pv.text = val
        pv.font.size = Pt(13)
        pv.font.bold = True
        pv.font.color.rgb = TEAL_ACCENT
        pv.space_before = Pt(4)

    # Right Before vs After Card
    c_bva = add_gray_card(slide9, Inches(8.3), Inches(2.1), Inches(4.233), Inches(3.0))
    tf_bva = c_bva.text_frame
    tf_bva.margin_left = tf_bva.margin_top = tf_bva.margin_right = tf_bva.margin_bottom = Inches(0.15)
    tf_bva.word_wrap = True
    pt_bva = tf_bva.paragraphs[0]
    pt_bva.text = "Before vs After comparison"
    pt_bva.font.size = Pt(14)
    pt_bva.font.bold = True
    pt_bva.font.color.rgb = NAVY_HEADER

    pv_bva = tf_bva.add_paragraph()
    pv_bva.text = "• Before: High cloud API costs, broken font encodings, trivial fill-in-blank MCQs.\n\n• After: 100% free local Jetson execution, regex text repair, 92% conceptual MCQ accuracy."
    pv_bva.font.size = Pt(10.5)
    pv_bva.font.color.rgb = TEXT_MUTED
    pv_bva.space_before = Pt(6)

    # Bottom Teal Box
    t_box9 = add_teal_card(slide9, Inches(8.3), Inches(5.3), Inches(4.233), Inches(1.2), title_text="What to include")
    tf_tb9 = t_box9.text_frame
    p_tb9_t = tf_tb9.paragraphs[0]
    p_tb9_t.alignment = PP_ALIGN.LEFT
    p_tb9_b = tf_tb9.add_paragraph()
    p_tb9_b.text = "Use actual numbers. Show only the proof points that matter most."
    p_tb9_b.font.size = Pt(11)
    p_tb9_b.font.color.rgb = TEXT_MAIN
    p_tb9_b.space_before = Pt(2)

    # =========================================================================
    # SLIDE 10: What Makes This Project Interesting, Conclusion & Next Steps
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_white_bg(slide10)
    add_header_and_footer(slide10, "What Makes This Project Interesting, Conclusion & Next Steps", "Wrap up with novelty, learning, and next steps.", 10)

    s10_cols = [
        ("What Makes It Interesting", "Solves the 1B SLM instruction-overload problem by delegating option shuffling and validation to Python while running regex text repair prior to vector inference."),
        ("Conclusion / Key Learning", "Demonstrates that lightweight 1B models deployed on edge hardware can outperform cloud APIs in reliability and domain-specific assessment tasks when paired with strict guardrails."),
        ("Next Steps", "1. Voice-guided quiz audio narration.\n2. Multimodal lecture diagram parsing.\n3. Quantized 3B parameter model support on Jetson Orin."),
    ]

    c10_w = Inches(3.7)
    c10_gap = Inches(0.3)

    for i, (title, desc) in enumerate(s10_cols):
        card = add_gray_card(slide10, Inches(0.8) + i * (c10_w + c10_gap), Inches(2.1), c10_w, Inches(3.1))
        tf = card.text_frame
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.2)
        tf.word_wrap = True

        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(14)
        pt.font.bold = True
        pt.font.color.rgb = NAVY_HEADER

        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(11)
        pd.font.color.rgb = TEXT_MUTED
        pd.space_before = Pt(8)

    # Bottom Left Teal Box (Thank You / Q&A)
    t_box10a = add_teal_card(slide10, Inches(0.8), Inches(5.4), Inches(7.5), Inches(1.1), title_text="Thank You / Q&A")
    tf_tb10a = t_box10a.text_frame
    p_tb10a_t = tf_tb10a.paragraphs[0]
    p_tb10a_t.alignment = PP_ALIGN.LEFT
    p_tb10a_b = tf_tb10a.add_paragraph()
    p_tb10a_b.text = "Thank you for your time. Happy to take questions and live feedback."
    p_tb10a_b.font.size = Pt(11)
    p_tb10a_b.font.color.rgb = TEXT_MAIN
    p_tb10a_b.space_before = Pt(2)

    # Bottom Right Teal Box (What to include)
    t_box10b = add_teal_card(slide10, Inches(8.533), Inches(5.4), Inches(4.0), Inches(1.1), title_text="What to include")
    tf_tb10b = t_box10b.text_frame
    p_tb10b_t = tf_tb10b.paragraphs[0]
    p_tb10b_t.alignment = PP_ALIGN.LEFT
    p_tb10b_b = tf_tb10b.add_paragraph()
    p_tb10b_b.text = "End with clarity, not clutter. Keep the last slide calm and confident."
    p_tb10b_b.font.size = Pt(10.5)
    p_tb10b_b.font.color.rgb = TEXT_MAIN
    p_tb10b_b.space_before = Pt(2)

    # =========================================================================
    # SLIDE 11: Project Evaluation Rubric
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_white_bg(slide11)
    add_header_and_footer(slide11, "Project Evaluation Rubric", "Each criterion will be scored from 1 to 10. Weighted total generates a score out of 100.", 11)

    # Top Formula Box
    t_formula = add_teal_card(slide11, Inches(0.8), Inches(1.9), Inches(11.733), Inches(0.7), title_text="")
    tf_tf = t_formula.text_frame
    tf_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_tf = tf_tf.paragraphs[0]
    p_tf.text = "Formula: Total Score = Σ [(Jury Rating ÷ 10) × Criterion Weight]    Example: Functionality rating 8/10 = 24/30"
    p_tf.font.size = Pt(11)
    p_tf.font.bold = True
    p_tf.font.color.rgb = TEXT_MAIN

    # Rubric Table
    rows, cols = 6, 5
    left, top, width, height = Inches(0.8), Inches(2.7), Inches(11.733), Inches(2.6)
    t_shape = slide11.shapes.add_table(rows, cols, left, top, width, height)
    t_rubric = t_shape.table

    t_rubric.columns[0].width = Inches(2.2)
    t_rubric.columns[1].width = Inches(0.9)
    t_rubric.columns[2].width = Inches(6.433)
    t_rubric.columns[3].width = Inches(1.1)
    t_rubric.columns[4].width = Inches(1.1)

    r_headers = ["Criterion", "Wt.", "What to evaluate", "Rating", "Score"]
    for j, h in enumerate(r_headers):
        cell = t_rubric.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY_HEADER
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

    r_data = [
        ["Functionality", "30", "Works on Jetson, user can complete the core task, output is useful and reliable.", "1–10", "__ /30"],
        ["Innovation", "25", "Novel idea, real problem fit, creative implementation, goes beyond starter template.", "1–10", "__ /25"],
        ["Edge-Readiness", "20", "Offline operation, latency, memory usage, efficient model/tool usage on Jetson.", "1–10", "__ /20"],
        ["Technical Quality", "15", "Clean architecture, readable code, error handling, README, meaningful Git history.", "1–10", "__ /15"],
        ["Presentation & Demo", "10", "Clear story, confident live demo, explains why edge deployment is the right choice.", "1–10", "__ /10"],
    ]

    for i, row in enumerate(r_data):
        for j, val in enumerate(row):
            cell = t_rubric.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = GRAY_CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(10)
            p.font.color.rgb = TEAL_ACCENT if j == 1 else (NAVY_HEADER if j == 0 else TEXT_MUTED)
            if j in [0, 1]:
                p.font.bold = True

    # Bottom Left Box (Rating anchors)
    c_r1 = add_gray_card(slide11, Inches(0.8), Inches(5.5), Inches(5.7), Inches(1.0))
    tf_r1 = c_r1.text_frame
    tf_r1.margin_left = tf_r1.margin_top = tf_r1.margin_right = tf_r1.margin_bottom = Inches(0.12)
    tf_r1.word_wrap = True
    pt_r1 = tf_r1.paragraphs[0]
    pt_r1.text = "Rating anchors"
    pt_r1.font.size = Pt(12)
    pt_r1.font.bold = True
    pt_r1.font.color.rgb = NAVY_HEADER

    pv_r1 = tf_r1.add_paragraph()
    pv_r1.text = "1–2 Poor  |  3–4 Basic  |  5–6 Working  |  7–8 Strong  |  9–10 Excellent"
    pv_r1.font.size = Pt(10.5)
    pv_r1.font.color.rgb = TEXT_MUTED
    pv_r1.space_before = Pt(2)

    # Bottom Right Box (Ranking rules)
    c_r2 = add_gray_card(slide11, Inches(6.833), Inches(5.5), Inches(5.7), Inches(1.0))
    tf_r2 = c_r2.text_frame
    tf_r2.margin_left = tf_r2.margin_top = tf_r2.margin_right = tf_r2.margin_bottom = Inches(0.12)
    tf_r2.word_wrap = True
    pt_r2 = tf_r2.paragraphs[0]
    pt_r2.text = "Ranking rules"
    pt_r2.font.size = Pt(12)
    pt_r2.font.bold = True
    pt_r2.font.color.rgb = NAVY_HEADER

    pv_r2 = tf_r2.add_paragraph()
    pv_r2.text = "Rank by total score out of 100. Tie-breakers: Functionality ➔ Edge-Readiness ➔ Innovation."
    pv_r2.font.size = Pt(10)
    pv_r2.font.color.rgb = TEXT_MUTED
    pv_r2.space_before = Pt(2)

    # Save presentation
    output_path = "edumind_final_presentation.pptx"
    prs.save(output_path)
    print(f"Successfully generated exact template match '{output_path}' with 11 slides.")


if __name__ == "__main__":
    create_exact_aiproff_deck()
